from __future__ import annotations

import csv
import io
import re
from pathlib import Path

import fitz
from openpyxl import load_workbook
from pptx import Presentation

from src.models import FileDiagnostic


TEXT_EXTENSIONS = {".txt", ".csv", ".tsv", ".md"}


def extension(filename: str) -> str:
    return Path(filename).suffix.lower()


def decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text(filename: str, data: bytes) -> tuple[str, dict[str, int | str]]:
    ext = extension(filename)
    if ext in TEXT_EXTENSIONS:
        text = decode_text(data)
        rows = None
        if ext in {".csv", ".tsv"}:
            dialect = "excel-tab" if ext == ".tsv" else "excel"
            try:
                rows = sum(1 for _ in csv.reader(io.StringIO(text), dialect=dialect))
            except csv.Error:
                rows = len(text.splitlines())
        return text, {"rows": rows or 0}
    if ext == ".pdf":
        document = fitz.open(stream=data, filetype="pdf")
        chunks = [page.get_text("text") for page in document]
        text = "\n".join(chunks)
        warning = ""
        if not text.strip():
            warning = "No embedded text found; this may be a scanned or image-only PDF."
        return text, {"pages": len(document), "warning": warning}
    if ext in {".pptx", ".ppt"}:
        if ext == ".ppt":
            raise ValueError("Legacy .ppt extraction is not supported; save it as .pptx.")
        presentation = Presentation(io.BytesIO(data))
        chunks: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks), {"slides": len(presentation.slides)}
    if ext in {".xlsx", ".xls"}:
        if ext == ".xls":
            raise ValueError("Legacy .xls extraction is not supported; save it as .xlsx or CSV.")
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        chunks: list[str] = []
        total_rows = 0
        for sheet in workbook.worksheets:
            chunks.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    chunks.append("\t".join(values))
                    total_rows += 1
        return "\n".join(chunks), {"sheets": len(workbook.sheetnames), "rows": total_rows}
    raise ValueError(f"Unsupported file extension: {ext or '(none)'}")


def diagnose(filename: str, data: bytes) -> FileDiagnostic:
    ext = extension(filename)
    try:
        text, metadata = extract_text(filename, data)
        return FileDiagnostic(
            filename=filename,
            extension=ext,
            bytes=len(data),
            characters=len(text),
            words=len(re.findall(r"\S+", text)),
            pages=_optional_int(metadata.get("pages")),
            slides=_optional_int(metadata.get("slides")),
            sheets=_optional_int(metadata.get("sheets")),
            rows=_optional_int(metadata.get("rows")),
            warning=str(metadata.get("warning", "")),
        )
    except Exception as exc:
        return FileDiagnostic(filename=filename, extension=ext, bytes=len(data), warning=str(exc))


def _optional_int(value: object) -> int | None:
    if value is None:
        return None
    number = int(value)
    return number if number > 0 else None


def mime_type(filename: str) -> str:
    return {
        ".pdf": "application/pdf",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".csv": "text/csv",
        ".tsv": "text/tab-separated-values",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }.get(extension(filename), "application/octet-stream")

