import io

from openpyxl import Workbook
from pptx import Presentation

from src.files import diagnose, extract_text


def test_text_diagnostics():
    result = diagnose("a.csv", b"name,value\nalpha,1\nbeta,2\n")
    assert result.rows == 3
    assert result.words >= 3
    assert result.characters > 0


def test_xlsx_extracts_sheets_and_rows():
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["name", "value"])
    sheet.append(["alpha", 1])
    buffer = io.BytesIO()
    workbook.save(buffer)
    text, metadata = extract_text("a.xlsx", buffer.getvalue())
    assert "Sheet: Data" in text
    assert "alpha" in text
    assert metadata["sheets"] == 1
    assert metadata["rows"] == 2


def test_pptx_extracts_slide_text():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Token costs"
    buffer = io.BytesIO()
    presentation.save(buffer)
    text, metadata = extract_text("a.pptx", buffer.getvalue())
    assert "Token costs" in text
    assert metadata["slides"] == 1
